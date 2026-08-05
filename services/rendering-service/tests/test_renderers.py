from io import BytesIO

import pytest
from docx import Document
from odf.opendocument import OpenDocumentSpreadsheet
from odf.table import Table, TableCell, TableRow
from odf.text import P
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
from rendering_service.renderers import _libreoffice, select_renderers
from rendering_service.renderers._libreoffice import ConversionError
from rendering_service.renderers.docx_text import DocxTextExtractionRenderer
from rendering_service.renderers.ods_text import OdsTextExtractionRenderer
from rendering_service.renderers.pdf_archive import PdfArchiveRenderer
from rendering_service.renderers.pptx_text import PptxTextExtractionRenderer
from rendering_service.renderers.thumbnail import ThumbnailRenderer
from reportlab.pdfgen import canvas


def _real_png(size=(800, 600)) -> bytes:
    buf = BytesIO()
    Image.new("RGB", size, color=(120, 20, 200)).save(buf, format="PNG")
    return buf.getvalue()


def _real_docx(text: str) -> bytes:
    doc = Document()
    doc.add_paragraph(text)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _real_pptx(title: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _real_ods(sheet_name: str, cell_text: str) -> bytes:
    doc = OpenDocumentSpreadsheet()
    table = Table(name=sheet_name)
    row = TableRow()
    cell = TableCell()
    cell.addElement(P(text=cell_text))
    row.addElement(cell)
    table.addElement(row)
    doc.spreadsheet.addElement(table)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _real_pdf(pages: int = 2) -> bytes:
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=(200, 200))
    for i in range(pages):
        c.drawString(10, 100, f"Seite {i + 1}")
        c.showPage()
    c.save()
    return buf.getvalue()


@pytest.mark.asyncio
async def test_thumbnail_renderer_downsizes_and_supports_image():
    renderer = ThumbnailRenderer()
    assert renderer.supports(content_type="image/png", filename="foo.png")
    assert renderer.supports(content_type=None, filename="foo.PNG")
    assert not renderer.supports(content_type="application/pdf", filename="foo.pdf")

    output = await renderer.render(_real_png(), filename="foto.png", content_type="image/png")
    assert output.target_content_type == "image/png"
    assert output.target_filename == "foto_thumbnail.png"
    thumb = Image.open(BytesIO(output.data))
    assert thumb.width <= 256
    assert thumb.height <= 256


@pytest.mark.asyncio
async def test_docx_text_extraction_renderer():
    renderer = DocxTextExtractionRenderer()
    assert renderer.supports(content_type=None, filename="brief.docx")
    assert not renderer.supports(content_type=None, filename="brief.txt")

    data = _real_docx("Hallo Welt, dies ist ein Test.")
    output = await renderer.render(data, filename="brief.docx", content_type=None)
    assert output.target_filename == "brief.txt"
    assert output.target_content_type == "text/plain; charset=utf-8"
    assert b"Hallo Welt, dies ist ein Test." in output.data


@pytest.mark.asyncio
async def test_pptx_text_extraction_renderer():
    renderer = PptxTextExtractionRenderer()
    assert renderer.supports(content_type=None, filename="folien.pptx")

    data = _real_pptx("Quartalsbericht")
    output = await renderer.render(data, filename="folien.pptx", content_type=None)
    assert output.target_filename == "folien.txt"
    assert b"Quartalsbericht" in output.data
    assert b"Folie 1" in output.data


@pytest.mark.asyncio
async def test_ods_text_extraction_renderer():
    renderer = OdsTextExtractionRenderer()
    assert renderer.supports(content_type=None, filename="tabelle.ods")
    assert renderer.supports(
        content_type="application/vnd.oasis.opendocument.spreadsheet", filename="a"
    )
    assert not renderer.supports(content_type=None, filename="tabelle.xlsx")

    data = _real_ods("Kosten", "Gesamtsumme")
    output = await renderer.render(data, filename="tabelle.ods", content_type=None)
    assert output.target_filename == "tabelle.txt"
    assert output.target_content_type == "text/plain; charset=utf-8"
    assert b"Kosten" in output.data
    assert b"Gesamtsumme" in output.data


@pytest.mark.asyncio
async def test_pdf_archive_renderer_preserves_pages():
    renderer = PdfArchiveRenderer()
    assert renderer.supports(content_type="application/pdf", filename="akte.pdf")

    data = _real_pdf(pages=3)
    output = await renderer.render(data, filename="akte.pdf", content_type="application/pdf")
    assert output.target_filename == "akte_archiv.pdf"
    assert output.target_content_type == "application/pdf"
    reader = PdfReader(BytesIO(output.data))
    assert len(reader.pages) == 3


@pytest.mark.asyncio
async def test_pdf_archive_renderer_converts_raster_image_via_pillow():
    """5.6 (seit P7-S3): Rasterbilder werden direkt ueber Pillow zu PDF
    gerendert, ohne den LibreOffice-Subprozess zu benoetigen."""
    renderer = PdfArchiveRenderer()
    assert renderer.supports(content_type="image/png", filename="foto.png")
    assert renderer.supports(content_type=None, filename="foto.jpg")

    output = await renderer.render(_real_png(), filename="foto.png", content_type="image/png")
    assert output.target_filename == "foto_archiv.pdf"
    assert output.target_content_type == "application/pdf"
    assert output.data.startswith(b"%PDF")
    reader = PdfReader(BytesIO(output.data))
    assert len(reader.pages) == 1


@pytest.mark.asyncio
async def test_pdf_archive_renderer_converts_docx_via_libreoffice():
    """5.6 (seit P7-S3): alle gaengigen Office-/Textformate muessen
    aussonderungsfaehig sein (Nutzervorgabe), nicht nur bereits-PDF-
    Dokumente - Konvertierung ueber LibreOffice headless."""
    renderer = PdfArchiveRenderer()
    assert renderer.supports(content_type=None, filename="brief.docx")

    data = _real_docx("Hallo Aussonderung, dies ist ein Test.")
    output = await renderer.render(data, filename="brief.docx", content_type=None)
    assert output.target_filename == "brief_archiv.pdf"
    assert output.target_content_type == "application/pdf"
    assert output.data.startswith(b"%PDF")


@pytest.mark.asyncio
async def test_pdf_archive_renderer_converts_plain_text_via_libreoffice():
    renderer = PdfArchiveRenderer()
    assert renderer.supports(content_type=None, filename="notiz.txt")

    output = await renderer.render(
        b"Einfacher Text zur Aussonderung.", filename="notiz.txt", content_type="text/plain"
    )
    assert output.data.startswith(b"%PDF")


@pytest.mark.asyncio
async def test_pdf_archive_renderer_raises_when_libreoffice_binary_missing(monkeypatch):
    """LibreOffice ist bemerkenswert nachsichtig beim Parsen (auch nicht
    valide .docx-Bytes werden meist noch als Text importiert) - der
    realistische Fehlerfall ist ein fehlendes Binary, nicht ein
    "unparsbares" Dokument."""
    monkeypatch.setattr(_libreoffice, "_BINARY_CANDIDATES", ("does-not-exist-binary",))
    renderer = PdfArchiveRenderer()
    with pytest.raises(ConversionError):
        await renderer.render(b"beliebiger Inhalt", filename="brief.docx", content_type=None)


def test_select_renderers_matches_expected_rules():
    # Seit P7-S3 (5.6) deckt PdfArchiveRenderer zusaetzlich Rasterbilder und
    # LibreOffice-konvertierbare Office-/Textformate ab (nicht mehr nur
    # bereits-PDF-Dokumente) - jede dieser Regeln liefert daher zusaetzlich
    # "pdf_archive" neben der jeweiligen formatspezifischen Regel.
    image_renderers = select_renderers(content_type="image/jpeg", filename="a.jpg")
    assert {r.rendition_type for r in image_renderers} == {"thumbnail", "pdf_archive"}

    docx_renderers = select_renderers(content_type=None, filename="a.docx")
    assert {r.rendition_type for r in docx_renderers} == {"substitute_text", "pdf_archive"}

    ods_renderers = select_renderers(content_type=None, filename="a.ods")
    assert {r.rendition_type for r in ods_renderers} == {"substitute_text", "pdf_archive"}

    csv_renderers = select_renderers(content_type="text/csv", filename="a.csv")
    assert {r.rendition_type for r in csv_renderers} == {"pdf_archive"}
