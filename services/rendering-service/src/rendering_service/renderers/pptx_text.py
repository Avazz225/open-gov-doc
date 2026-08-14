import os
from io import BytesIO

from pptx import Presentation

from rendering_service.renderers.interface import Renderer, RenderOutput

_PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class PptxTextExtractionRenderer(Renderer):
    """Rendition for PowerPoint presentations (2.4): text extraction to .txt
    instead of the .pdf conversion named as an example in the concept -
    a real .pptx->.pdf conversion would need an external Office rendering
    component (e.g. LibreOffice headless), which is not reliably/quickly
    available in this environment (the same trade-off as ClamAV vs. the
    EICAR engine in P5-S1, see ADR 0010 and PROGRESS.md). Text nonetheless
    remains a fully valid, fail-safe rendition in the sense of 2.4."""

    rendition_type = "substitute_text"

    def supports(self, *, content_type: str | None, filename: str) -> bool:
        if content_type == _PPTX_CONTENT_TYPE:
            return True
        return filename.lower().endswith(".pptx")

    async def render(self, data: bytes, *, filename: str, content_type: str | None) -> RenderOutput:
        presentation = Presentation(BytesIO(data))
        slide_texts = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text
            ]
            slide_texts.append(f"--- Folie {index} ---\n" + "\n".join(texts))
        stem = os.path.splitext(filename)[0] or "praesentation"
        return RenderOutput(
            rendition_type=self.rendition_type,
            target_filename=f"{stem}.txt",
            target_content_type="text/plain; charset=utf-8",
            data="\n\n".join(slide_texts).encode("utf-8"),
        )
